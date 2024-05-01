from flask import Flask, render_template, request, send_file
from PyPDF2 import PdfReader
from pptx import Presentation
import os

app = Flask(__name__)

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if request.method == 'POST':
        file = request.files['file']
        if file:
            pdf_text = ""
            pdf_reader = PdfReader(file)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                pdf_text += page.extract_text()

            prs = Presentation()
            slide_layout = prs.slide_layouts[1] 

            for text in pdf_text.split('\n\n'):  
                slide = prs.slides.add_slide(slide_layout)
                title, content = text.strip().split('\n', 1) if '\n' in text else (text.strip(), '')
                
                title_shape = None
                content_placeholder = None
                for shape in slide.placeholders:
                    if shape.is_placeholder:
                        if shape.placeholder_format.idx == 0:
                            title_shape = shape
                        elif shape.placeholder_format.idx == 1:
                            content_placeholder = shape
                
                if title_shape:
                    title_shape.text = title
                
                if content_placeholder:
                    content_placeholder.text = content

            output_filename = "output.pptx"
            prs.save(output_filename)

            return render_template('index.html', alert_message="Conversion successful!", output_filename=output_filename)

    return render_template('index.html', alert_message="No file selected!")

@app.route('/download/<filename>', methods=['GET', 'POST'])
def download(filename):
    if os.path.exists(filename):
        return send_file(filename, as_attachment=True)
    else:
        return "File not found"

if __name__ == '__main__':
    app.run(debug=True)
